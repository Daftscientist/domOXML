"""PresentationML reader tests against deterministic domOXML-generated decks."""

# pyright: reportPrivateUsage=false

from __future__ import annotations

from io import BytesIO
from xml.etree.ElementTree import Element

import pytest
from defusedxml import ElementTree
from PIL import Image, ImageChops
from pptx import Presentation as PptxPresentation
from pptx.util import Inches

from domoxml import Presentation, pptx_to_html
from domoxml.core.ir.model import (
    AutoNumberBullet,
    Blur,
    Box,
    CharBullet,
    FillOverlay,
    GradientFill,
    GradientStop,
    GroupNode,
    Line,
    PictureFill,
    PortableFallback,
    PreservedNode,
    Reflection,
    Rgba,
    Shadow,
    ShapeNode,
    SlideIR,
    SoftEdge,
    SolidFill,
    SourceProvenance,
    TextBody,
    TextParagraph,
    TextRun,
    Transform,
)
from domoxml.core.opc import OpcPackage, write_package
from domoxml.core.roundtrip import render_html_roundtrip
from domoxml.slides import build_pptx, read_pptx, read_pptx_result
from domoxml.slides.appearance_read import rgba
from domoxml.slides.read import (
    _can_own_source_shape_crop,
    _connector_reverse_coverage,
    _group_reverse_coverage,
    _slide_colors,
    _with_pptx_identity,
)
from domoxml.types import Editability, HtmlPresentation, Representation, SourceRetention

_A = "http://schemas.openxmlformats.org/drawingml/2006/main"
_MC = "http://schemas.openxmlformats.org/markup-compatibility/2006"
_P = "http://schemas.openxmlformats.org/presentationml/2006/main"
_R = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"
_PKG_REL = "http://schemas.openxmlformats.org/package/2006/relationships"


def _sample_ir() -> SlideIR:
    return SlideIR(
        width=12_192_000,
        height=6_858_000,
        shapes=(
            ShapeNode(
                box=Box(x=914_400, y=914_400, width=3_657_600, height=1_828_800),
                geom="roundRect",
                fill=SolidFill(color=Rgba(r=79, g=70, b=229, a=0.75)),
                line=Line(color=Rgba(r=1, g=2, b=3), width_emu=19_050, dash="dash"),
                effects=(
                    Shadow(
                        color=Rgba(r=0, g=0, b=0, a=0.4),
                        blur_emu=20_000,
                        distance_emu=10_000,
                        direction_deg=90,
                    ),
                ),
                corner_radius_emu=76_200,
                text=TextBody(
                    paragraphs=(
                        TextParagraph(
                            runs=(
                                TextRun(text="Coffee ", font_family="Inter", size_pt=24),
                                TextRun(
                                    text="calm",
                                    font_family="Inter",
                                    size_pt=24,
                                    italic=True,
                                    color=Rgba(r=255, g=255, b=255),
                                ),
                            ),
                            align="center",
                        ),
                    )
                ),
            ),
            ShapeNode(
                box=Box(x=0, y=0, width=100, height=100),
                fill=GradientFill(
                    stops=(
                        GradientStop(pos=0, color=Rgba(r=0, g=0, b=0)),
                        GradientStop(pos=1, color=Rgba(r=255, g=255, b=255)),
                    ),
                    angle_deg=90,
                ),
            ),
            ShapeNode(
                box=Box(x=100, y=100, width=200, height=200),
                fill=PictureFill(data=b"png-bytes", ext="png"),
            ),
        ),
    )


def test_reads_generated_pptx_into_canvas_ir() -> None:
    pptx = build_pptx([_sample_ir()], faces=[])
    result = read_pptx_result(pptx)
    [slide] = result.slides
    assert (slide.width, slide.height) == (12_192_000, 6_858_000)
    assert len(slide.shapes) == 3

    rich, gradient, picture = slide.shapes
    assert rich.box.x == 914_400 and rich.geom == "roundRect"
    assert isinstance(rich.fill, SolidFill) and rich.fill.color.a == 0.75
    assert rich.line is not None and rich.line.dash == "dash"
    assert rich.shadow is not None and rich.shadow.direction_deg == 90
    assert rich.text is not None
    assert "".join(run.text for run in rich.text.paragraphs[0].runs) == "Coffee calm"
    assert rich.text.paragraphs[0].runs[1].italic is True

    assert isinstance(gradient.fill, GradientFill)
    assert gradient.fill.angle_deg == 90
    assert isinstance(picture.fill, PictureFill)
    assert picture.fill.data == b"png-bytes"
    assert result.coverage.count(Representation.NATIVE) == 3
    assert result.coverage.count_editability(Editability.SEMANTIC) == 3
    assert result.coverage.output_count == 3
    assert result.coverage.raster_area_emu2 == 0


def test_portable_blur_fallback_uses_alternate_content_and_round_trips() -> None:
    fallback_box = Box(x=800_000, y=700_000, width=2_400_000, height=1_400_000)
    shape = ShapeNode(
        box=Box(x=1_000_000, y=900_000, width=2_000_000, height=1_000_000),
        fill=SolidFill(color=Rgba(r=232, g=74, b=95)),
        effects=(Blur(radius_emu=95_250),),
        portable_fallback=PortableFallback(
            box=fallback_box,
            picture=PictureFill(
                data=b"isolated-blur-png",
                ext="png",
                raster_role="portable-blur-fallback",
            ),
        ),
    )

    pptx = build_pptx([SlideIR(width=12_192_000, height=6_858_000, contents=(shape,))], faces=[])
    slide_xml = OpcPackage.from_bytes(pptx).read("ppt/slides/slide1.xml").decode()

    assert 'mc:Choice Requires="p16"' in slide_xml
    assert '<a:blur rad="95250" grow="1"/>' in slide_xml
    assert 'descr="domoxml-raster:portable-blur-fallback"' in slide_xml
    assert '<mc:Choice Requires="p16"><p:sp>' in slide_xml
    assert slide_xml.count("domoxml-raster:portable-blur-fallback") == 2

    result = read_pptx_result(pptx)
    [recovered] = result.slides[0].shapes
    assert recovered.effects == shape.effects
    assert recovered.portable_fallback is not None
    assert recovered.portable_fallback.box == fallback_box
    assert recovered.portable_fallback.picture.data == b"isolated-blur-png"
    assert result.coverage.count(Representation.HYBRID) == 1
    assert result.coverage.count_editability(Editability.COMPONENTS) == 1
    assert result.coverage.output_count == 2
    assert result.coverage.raster_area_emu2 == fallback_box.width * fallback_box.height


def test_portable_reflection_fallback_uses_alternate_content_and_round_trips() -> None:
    fallback_box = Box(x=900_000, y=700_000, width=2_200_000, height=2_100_000)
    reflection = Reflection(
        distance_emu=114_300,
        start_alpha=0.8,
        end_alpha=0.0,
    )
    shape = ShapeNode(
        box=Box(x=1_000_000, y=900_000, width=2_000_000, height=900_000),
        fill=SolidFill(color=Rgba(r=232, g=74, b=95)),
        effects=(reflection,),
        portable_fallback=PortableFallback(
            box=fallback_box,
            picture=PictureFill(
                data=b"isolated-reflection-png",
                ext="png",
                raster_role="portable-effect-fallback",
            ),
        ),
    )

    pptx = build_pptx([SlideIR(width=12_192_000, height=6_858_000, contents=(shape,))], faces=[])
    slide_xml = OpcPackage.from_bytes(pptx).read("ppt/slides/slide1.xml").decode()

    assert 'mc:Choice Requires="p16"' in slide_xml
    assert '<a:reflection blurRad="0" dist="114300" stA="80000" endA="0"' in slide_xml
    assert slide_xml.count("domoxml-raster:portable-effect-fallback") == 2

    result = read_pptx_result(pptx)
    [recovered] = result.slides[0].shapes
    assert recovered.effects == (reflection,)
    assert recovered.portable_fallback is not None
    assert recovered.portable_fallback.box == fallback_box
    assert recovered.portable_fallback.picture.data == b"isolated-reflection-png"
    assert result.coverage.count(Representation.HYBRID) == 1
    assert result.coverage.count_editability(Editability.COMPONENTS) == 1
    assert result.coverage.output_count == 2
    assert result.coverage.raster_area_emu2 == fallback_box.width * fallback_box.height


def test_portable_soft_edge_fallback_uses_alternate_content_and_round_trips() -> None:
    fallback_box = Box(x=1_000_000, y=900_000, width=2_000_000, height=1_000_000)
    soft_edge = SoftEdge(radius_emu=114_300)
    shape = ShapeNode(
        box=fallback_box,
        fill=SolidFill(color=Rgba(r=39, g=102, b=120)),
        effects=(soft_edge,),
        portable_fallback=PortableFallback(
            box=fallback_box,
            picture=PictureFill(
                data=b"isolated-soft-edge-png",
                ext="png",
                raster_role="portable-effect-fallback",
            ),
        ),
    )

    pptx = build_pptx([SlideIR(width=12_192_000, height=6_858_000, contents=(shape,))], faces=[])
    slide_xml = OpcPackage.from_bytes(pptx).read("ppt/slides/slide1.xml").decode()

    assert 'mc:Choice Requires="p16"' in slide_xml
    assert '<a:softEdge rad="114300"/>' in slide_xml
    assert slide_xml.count("domoxml-raster:portable-effect-fallback") == 2

    result = read_pptx_result(pptx)
    [recovered] = result.slides[0].shapes
    assert recovered.effects == (soft_edge,)
    assert recovered.portable_fallback is not None
    assert recovered.portable_fallback.box == fallback_box
    assert recovered.portable_fallback.picture.data == b"isolated-soft-edge-png"
    assert result.coverage.count(Representation.HYBRID) == 1
    assert result.coverage.count_editability(Editability.COMPONENTS) == 1
    assert result.coverage.output_count == 2
    assert result.coverage.raster_area_emu2 == fallback_box.width * fallback_box.height


def test_portable_fill_overlay_fallback_uses_alternate_content_and_round_trips() -> None:
    fallback_box = Box(x=1_000_000, y=900_000, width=2_000_000, height=1_000_000)
    overlay = FillOverlay(
        fill=SolidFill(color=Rgba(r=255, g=40, b=80, a=0.75)),
        blend="mult",
    )
    shape = ShapeNode(
        box=fallback_box,
        fill=SolidFill(color=Rgba(r=20, g=60, b=140)),
        effects=(overlay,),
        portable_fallback=PortableFallback(
            box=fallback_box,
            picture=PictureFill(
                data=b"isolated-fill-overlay-png",
                ext="png",
                raster_role="portable-effect-fallback",
            ),
        ),
    )

    pptx = build_pptx([SlideIR(width=12_192_000, height=6_858_000, contents=(shape,))], faces=[])
    slide_xml = OpcPackage.from_bytes(pptx).read("ppt/slides/slide1.xml").decode()

    assert 'mc:Choice Requires="p16"' in slide_xml
    assert '<a:fillOverlay blend="mult"><a:solidFill>' in slide_xml
    assert '<a:srgbClr val="FF2850"><a:alpha val="75000"/>' in slide_xml
    assert slide_xml.count("domoxml-raster:portable-effect-fallback") == 1
    assert '<mc:Choice Requires="p16"><p:sp>' in slide_xml
    assert "</p:sp></mc:Choice>" in slide_xml

    result = read_pptx_result(pptx)
    [recovered] = result.slides[0].shapes
    assert recovered.effects == (overlay,)
    assert recovered.portable_fallback is not None
    assert recovered.portable_fallback.box == fallback_box
    assert recovered.portable_fallback.picture.data == b"isolated-fill-overlay-png"
    assert result.coverage.count(Representation.HYBRID) == 1
    assert result.coverage.count_editability(Editability.COMPONENTS) == 1
    assert result.coverage.output_count == 2
    assert result.coverage.raster_area_emu2 == fallback_box.width * fallback_box.height


def test_portable_fallback_reports_and_retains_unsupported_choice_effects() -> None:
    shape = ShapeNode(
        box=Box(x=1_000_000, y=900_000, width=2_000_000, height=1_000_000),
        fill=SolidFill(color=Rgba(r=232, g=74, b=95)),
        effects=(Blur(radius_emu=95_250),),
        portable_fallback=PortableFallback(
            box=Box(x=800_000, y=700_000, width=2_400_000, height=1_400_000),
            picture=PictureFill(data=b"isolated-blur-png", ext="png"),
        ),
    )
    package = OpcPackage.from_bytes(
        build_pptx([SlideIR(width=12_192_000, height=6_858_000, contents=(shape,))], faces=[])
    )
    parts: dict[str, bytes | str] = {part: package.read(part) for part in package.parts}
    slide_part = "ppt/slides/slide1.xml"
    root = ElementTree.fromstring(package.read(slide_part))
    effect_list = root.find(
        ".//{http://schemas.openxmlformats.org/markup-compatibility/2006}Choice/"
        "p:sp/p:spPr/a:effectLst",
        {"p": _P, "a": _A},
    )
    assert effect_list is not None
    effect_list.append(ElementTree.fromstring(f'<a:fillOverlay xmlns:a="{_A}" blend="mult"/>'))
    parts[slide_part] = ElementTree.tostring(root)

    result = read_pptx_result(write_package(parts))

    [recovered] = [node for node in result.slides[0].contents if isinstance(node, PreservedNode)]
    assert recovered.fallback is not None
    assert recovered.fallback.data == b"isolated-blur-png"
    assert shape.portable_fallback is not None
    assert recovered.box == shape.portable_fallback.box
    assert "fillOverlay" in recovered.payload.root_xml
    [coverage] = result.coverage.items
    assert coverage.representation is Representation.ELEMENT_LAYER
    assert coverage.editability is Editability.LAYERS
    assert coverage.source_retention is SourceRetention.ATTACHED
    assert coverage.raster_area_emu2 == (
        shape.portable_fallback.box.width * shape.portable_fallback.box.height
    )
    assert "owned visual layer" in (coverage.reason or "")
    assert {fragment.kind for fragment in result.preserved} == {"AlternateContent"}

    rebuilt = build_pptx(list(result.slides), faces=[])
    assert b"fillOverlay" in OpcPackage.from_bytes(rebuilt).read(slide_part)


def _preset_shadow_source(*, with_sibling: bool = False) -> bytes:
    shape = ShapeNode(
        box=Box(x=3_000_000, y=2_000_000, width=6_000_000, height=2_500_000),
        fill=SolidFill(color=Rgba(r=37, g=99, b=235)),
        effects=(
            Shadow(
                color=Rgba(r=20, g=20, b=25, a=0.65),
                blur_emu=0,
                distance_emu=550_000,
                direction_deg=45,
            ),
        ),
    )
    sibling = ShapeNode(
        box=Box(x=7_000_000, y=3_000_000, width=3_000_000, height=1_500_000),
        fill=SolidFill(color=Rgba(r=239, g=68, b=68, a=0.55)),
    )
    contents = (shape, sibling) if with_sibling else (shape,)
    package = OpcPackage.from_bytes(
        build_pptx([SlideIR(width=12_192_000, height=6_858_000, contents=contents)], faces=[])
    )
    slide_part = "ppt/slides/slide1.xml"
    parts: dict[str, bytes | str] = {part: package.read(part) for part in package.parts}
    root = ElementTree.fromstring(package.read(slide_part))
    effect_list = root.find(".//p:sp/p:spPr/a:effectLst", {"p": _P, "a": _A})
    assert effect_list is not None
    outer_shadow = effect_list.find("a:outerShdw", {"a": _A})
    assert outer_shadow is not None
    preset_shadow = Element(
        f"{{{_A}}}prstShdw",
        {"prst": "shdw3", "dist": "550000", "dir": "2700000"},
    )
    preset_shadow.extend(tuple(outer_shadow))
    effect_list.insert(list(effect_list).index(outer_shadow), preset_shadow)
    effect_list.remove(outer_shadow)
    parts[slide_part] = ElementTree.tostring(root)
    return write_package(parts)


def _append_preset_shadow_sibling(pptx: bytes) -> bytes:
    package = OpcPackage.from_bytes(pptx)
    sibling_package = OpcPackage.from_bytes(_preset_shadow_source(with_sibling=True))
    slide_part = "ppt/slides/slide1.xml"
    root = ElementTree.fromstring(package.read(slide_part))
    sibling_root = ElementTree.fromstring(sibling_package.read(slide_part))
    tree = root.find("p:cSld/p:spTree", {"p": _P})
    sibling_shapes = sibling_root.findall("p:cSld/p:spTree/p:sp", {"p": _P})
    assert tree is not None
    assert len(sibling_shapes) == 2
    tree.append(sibling_shapes[1])
    parts: dict[str, bytes | str] = {part: package.read(part) for part in package.parts}
    parts[slide_part] = ElementTree.tostring(root)
    return write_package(parts)


def test_preset_shadow_uses_full_slide_rasterized_fallback_and_re_emits() -> None:
    rendered = BytesIO()
    Image.new("RGB", (1280, 720), "#5D7893").save(rendered, "PNG")
    source = _preset_shadow_source()

    result = read_pptx_result(source, fallback_pngs=(rendered.getvalue(),))

    [preserved] = [node for node in result.slides[0].contents if isinstance(node, PreservedNode)]
    assert preserved.box == Box(x=0, y=0, width=12_192_000, height=6_858_000)
    assert preserved.fallback is not None
    assert preserved.fallback_representation == "rasterized"
    assert 'prst="shdw3"' in preserved.payload.root_xml
    [coverage] = result.coverage.items
    assert coverage.representation is Representation.RASTERIZED
    assert coverage.editability is Editability.NONE
    assert coverage.source_retention is SourceRetention.ATTACHED
    assert coverage.raster_area_emu2 == 12_192_000 * 6_858_000

    rebuilt = build_pptx(list(result.slides), faces=[])
    rebuilt_slide = OpcPackage.from_bytes(rebuilt).read("ppt/slides/slide1.xml")
    assert b'prstShdw prst="shdw3"' in rebuilt_slide
    assert b"AlternateContent" in rebuilt_slide
    assert b"domoxml-raster:pptx-source-rasterized" in rebuilt_slide

    recovered = read_pptx_result(rebuilt)
    [recovered_node] = [
        node for node in recovered.slides[0].contents if isinstance(node, PreservedNode)
    ]
    assert recovered_node.fallback is not None
    assert preserved.fallback is not None
    assert recovered_node.fallback.data == preserved.fallback.data
    assert recovered_node.fallback_representation == "rasterized"
    [recovered_coverage] = recovered.coverage.items
    assert recovered_coverage.representation is Representation.RASTERIZED
    assert recovered_coverage.editability is Editability.NONE
    assert recovered_coverage.source_retention is SourceRetention.ATTACHED
    assert recovered_coverage.raster_area_emu2 == 12_192_000 * 6_858_000


def test_preset_shadow_multi_visual_uses_slide_level_fallback() -> None:
    rendered = BytesIO()
    Image.new("RGB", (1280, 720), "#5D7893").save(rendered, "PNG")

    result = read_pptx_result(
        _preset_shadow_source(with_sibling=True),
        fallback_pngs=(rendered.getvalue(),),
    )

    slide = result.slides[0]
    assert slide.renderer_fallback is not None
    fallback_bytes = slide.renderer_fallback.data
    [preserved] = [node for node in slide.contents if isinstance(node, PreservedNode)]
    assert preserved.fallback is None
    assert 'prst="shdw3"' in preserved.payload.root_xml
    assert [item.representation for item in result.coverage.items] == [
        Representation.NATIVE,
        Representation.RASTERIZED,
    ]
    assert result.coverage.items[-1].source_retention is SourceRetention.ATTACHED
    assert result.coverage.items[-1].raster_area_emu2 == 12_192_000 * 6_858_000

    rebuilt = build_pptx(list(result.slides), faces=[])
    rebuilt_slide = OpcPackage.from_bytes(rebuilt).read("ppt/slides/slide1.xml")
    assert b'prstShdw prst="shdw3"' in rebuilt_slide
    assert rebuilt_slide.count(b"AlternateContent") == 2
    assert b"domoxml-raster:pptx-slide-rasterized" in rebuilt_slide

    recovered = read_pptx_result(rebuilt)
    recovered_slide = recovered.slides[0]
    assert recovered_slide.renderer_fallback is not None
    assert recovered_slide.renderer_fallback.data == fallback_bytes
    assert len(recovered_slide.contents) == 2
    assert [item.representation for item in recovered.coverage.items] == [
        Representation.NATIVE,
        Representation.RASTERIZED,
    ]


def test_malformed_slide_fallback_marker_stays_source_owned() -> None:
    rendered = BytesIO()
    Image.new("RGB", (1280, 720), "#5D7893").save(rendered, "PNG")
    result = read_pptx_result(
        _preset_shadow_source(with_sibling=True),
        fallback_pngs=(rendered.getvalue(),),
    )
    rebuilt = build_pptx(list(result.slides), faces=[])
    package = OpcPackage.from_bytes(rebuilt)
    slide_part = "ppt/slides/slide1.xml"
    root = ElementTree.fromstring(package.read(slide_part))
    offset = root.find(
        ".//mc:Fallback/p:pic/p:spPr/a:xfrm/a:off",
        {"mc": _MC, "p": _P, "a": _A},
    )
    assert offset is not None
    offset.set("x", "1")
    parts: dict[str, bytes | str] = {part: package.read(part) for part in package.parts}
    parts[slide_part] = ElementTree.tostring(root)

    recovered = read_pptx_result(write_package(parts))

    assert recovered.slides[0].renderer_fallback is None
    [preserved] = [node for node in recovered.slides[0].contents if isinstance(node, PreservedNode)]
    assert preserved.payload.kind == "AlternateContent"
    assert preserved.fallback_representation == "rasterized"
    [coverage] = recovered.coverage.items
    assert coverage.representation is Representation.RASTERIZED
    assert coverage.source_retention is SourceRetention.ATTACHED


def test_legacy_node_level_full_slide_fallback_rejects_siblings() -> None:
    rendered = BytesIO()
    Image.new("RGB", (1280, 720), "#5D7893").save(rendered, "PNG")

    safe_result = read_pptx_result(
        _preset_shadow_source(),
        fallback_pngs=(rendered.getvalue(),),
    )
    [fallback] = [
        node for node in safe_result.slides[0].contents if isinstance(node, PreservedNode)
    ]
    unsafe_rebuilt = _append_preset_shadow_sibling(build_pptx(list(safe_result.slides), faces=[]))
    unsafe_recovered = read_pptx_result(unsafe_rebuilt)
    assert not any(
        isinstance(node, PreservedNode) and node.fallback_representation == "rasterized"
        for node in unsafe_recovered.slides[0].contents
    )
    assert [item.representation for item in unsafe_recovered.coverage.items] == [
        Representation.APPROXIMATED,
        Representation.NATIVE,
    ]

    sibling = ShapeNode(
        box=Box(x=7_000_000, y=3_000_000, width=3_000_000, height=1_500_000),
        fill=SolidFill(color=Rgba(r=239, g=68, b=68, a=0.55)),
    )
    edited = safe_result.slides[0].model_copy(update={"contents": (fallback, sibling)})
    rebuilt_slide = OpcPackage.from_bytes(build_pptx([edited], faces=[])).read(
        "ppt/slides/slide1.xml"
    )
    assert b'prstShdw prst="shdw3"' in rebuilt_slide
    assert b"pptx-source-rasterized" not in rebuilt_slide


@pytest.mark.integration
def test_unsupported_over_overlay_uses_owned_source_render_and_re_emits() -> None:
    shape = ShapeNode(
        box=Box(x=1_000_000, y=900_000, width=2_000_000, height=1_000_000),
        fill=SolidFill(color=Rgba(r=20, g=60, b=140)),
        effects=(
            FillOverlay(
                fill=SolidFill(color=Rgba(r=255, g=40, b=80, a=0.75)),
                blend="mult",
            ),
        ),
        transform=Transform(rotation_deg=30),
    )
    package = OpcPackage.from_bytes(
        build_pptx([SlideIR(width=12_192_000, height=6_858_000, contents=(shape,))], faces=[])
    )
    parts: dict[str, bytes | str] = {part: package.read(part) for part in package.parts}
    slide_part = "ppt/slides/slide1.xml"
    parts[slide_part] = package.read(slide_part).replace(b'blend="mult"', b'blend="over"')
    rendered = BytesIO()
    Image.new("RGB", (1280, 720), "#5D7893").save(rendered, "PNG")

    current = write_package(parts)
    fallback_pngs = (rendered.getvalue(),)
    previous_metrics: tuple[int, int] | None = None
    previous_pngs: tuple[bytes, ...] | None = None
    for _cycle in range(2):
        html = Presentation.from_pptx(current, fallback_pngs=fallback_pngs)

        assert 'data-domoxml-representation="element-layer"' in html.slides[0].html
        assert "data-domoxml-preserved-payload=" in html.slides[0].html
        [fallback_asset] = [asset for asset in html.assets if asset.path.endswith(".png")]
        with Image.open(BytesIO(fallback_asset.data)) as fallback_image:
            alpha = fallback_image.convert("RGBA")
            corner = alpha.getpixel((0, 0))
            center = alpha.getpixel((alpha.width // 2, alpha.height // 2))
            assert isinstance(corner, tuple)
            assert isinstance(center, tuple)
            assert corner[3] == 0
            assert center[3] == 255
        [coverage] = html.coverage.items
        assert coverage.representation is Representation.ELEMENT_LAYER
        assert coverage.editability is Editability.LAYERS
        assert coverage.source_retention is SourceRetention.ATTACHED
        assert coverage.raster_area_emu2 > shape.box.width * shape.box.height
        metrics = (coverage.output_count, coverage.raster_area_emu2)
        if previous_metrics is not None:
            assert metrics == previous_metrics
        previous_metrics = metrics

        rebuilt = render_html_roundtrip(html)
        assert rebuilt.pptx is not None
        if previous_pngs is not None:
            with (
                Image.open(BytesIO(rebuilt.pngs[0])) as current_image,
                Image.open(BytesIO(previous_pngs[0])) as previous_image,
            ):
                assert ImageChops.difference(current_image, previous_image).getbbox() is None
        previous_pngs = rebuilt.pngs
        fallback_pngs = rebuilt.pngs
        current = rebuilt.pptx
        rebuilt_slide = OpcPackage.from_bytes(current).read(slide_part)
        assert b'fillOverlay blend="over"' in rebuilt_slide
        assert b"AlternateContent" in rebuilt_slide


def test_owned_source_crop_excludes_rounded_geometry() -> None:
    shape = ShapeNode(
        box=Box(x=0, y=0, width=2_000_000, height=1_000_000),
        fill=SolidFill(color=Rgba(r=20, g=60, b=140)),
    )

    assert _can_own_source_shape_crop(shape, is_only_visual=True)
    assert not _can_own_source_shape_crop(
        shape.model_copy(update={"corner_radius_emu": 100_000}),
        is_only_visual=True,
    )


def _overlapping_unsupported_overlay_html() -> tuple[HtmlPresentation, str]:
    overlay_shape = ShapeNode(
        box=Box(x=1_000_000, y=900_000, width=2_000_000, height=1_000_000),
        fill=SolidFill(color=Rgba(r=20, g=60, b=140)),
        effects=(
            FillOverlay(
                fill=SolidFill(color=Rgba(r=255, g=40, b=80, a=0.75)),
                blend="mult",
            ),
        ),
        transform=Transform(rotation_deg=30),
    )
    foreground = ShapeNode(
        box=Box(x=2_000_000, y=1_000_000, width=1_000_000, height=700_000),
        fill=SolidFill(color=Rgba(r=40, g=180, b=100)),
    )
    package = OpcPackage.from_bytes(
        build_pptx(
            [
                SlideIR(
                    width=12_192_000,
                    height=6_858_000,
                    contents=(overlay_shape, foreground),
                )
            ],
            faces=[],
        )
    )
    slide_part = "ppt/slides/slide1.xml"
    parts: dict[str, bytes | str] = {part: package.read(part) for part in package.parts}
    parts[slide_part] = package.read(slide_part).replace(b'blend="mult"', b'blend="over"', 1)
    rendered = BytesIO()
    Image.new("RGB", (1280, 720), "#5D7893").save(rendered, "PNG")

    return (
        Presentation.from_pptx(write_package(parts), fallback_pngs=(rendered.getvalue(),)),
        slide_part,
    )


def test_overlapping_unsupported_overlay_is_visible_without_false_layer_ownership() -> None:
    html, _ = _overlapping_unsupported_overlay_html()

    assert 'data-domoxml-representation="rasterized"' in html.slides[0].html
    assert "data-domoxml-preserved-payload=" in html.slides[0].html
    [coverage] = [
        item for item in html.coverage.items if item.representation is Representation.RASTERIZED
    ]
    assert coverage.editability is Editability.NONE
    assert coverage.source_retention is SourceRetention.ATTACHED
    assert "cannot prove independent ownership" in coverage.reason


@pytest.mark.integration
def test_overlapping_unsupported_overlay_re_emits_native_xml_without_portable_fallback() -> None:
    html, slide_part = _overlapping_unsupported_overlay_html()

    rebuilt = render_html_roundtrip(html)
    assert rebuilt.pptx is not None
    rebuilt_slide = OpcPackage.from_bytes(rebuilt.pptx).read(slide_part)
    assert b'fillOverlay blend="over"' in rebuilt_slide
    assert b"AlternateContent" not in rebuilt_slide


def test_exposes_generated_pptx_as_html() -> None:
    pptx = build_pptx([_sample_ir()], faces=[])
    html = pptx_to_html(pptx)
    assert len(html.slides) == 1
    assert "Coffee " in html.slides[0].html and "calm" in html.slides[0].html
    assert len(html.assets) == 1
    assert html.coverage.count(Representation.NATIVE) == 3
    assert Presentation.from_pptx(pptx) == html


def test_node_identity_and_provenance_round_trip_through_private_extension() -> None:
    node = ShapeNode(
        node_id="hero-title",
        provenance=SourceProvenance(
            source_format="html",
            source_id="title",
            owner_node_id="hero",
            role="title",
        ),
        box=Box(x=10, y=20, width=300, height=100),
        fill=SolidFill(color=Rgba(r=10, g=20, b=30)),
    )
    pptx = build_pptx([SlideIR(width=1_000, height=1_000, contents=(node,))], faces=[])

    slide_xml = OpcPackage.from_bytes(pptx).read("ppt/slides/slide1.xml").decode()
    assert 'dx:node xmlns:dx="urn:domoxml:canvas-ir:1" id="hero-title"' in slide_xml
    [recovered] = read_pptx(pptx)[0].contents

    assert recovered.node_id == "hero-title"
    assert recovered.provenance == node.provenance


def test_third_party_shape_uses_cnvpr_id_as_pptx_provenance() -> None:
    pptx = build_pptx([_sample_ir()], faces=[])
    # Remove domOXML's extension to exercise an ordinary producer's cNvPr fallback.
    package = OpcPackage.from_bytes(pptx)
    root = ElementTree.fromstring(package.read("ppt/slides/slide1.xml"))
    for nv_pr in root.findall(".//p:nvPr", {"p": _P}):
        ext_lst = nv_pr.find("p:extLst", {"p": _P})
        if ext_lst is not None:
            nv_pr.remove(ext_lst)
    parts: dict[str, bytes | str] = {part: package.read(part) for part in package.parts}
    parts["ppt/slides/slide1.xml"] = ElementTree.tostring(root)
    [shape] = read_pptx(write_package(parts))[0].contents[:1]

    assert shape.node_id == "pptx-2"
    assert shape.provenance == SourceProvenance(
        source_format="pptx",
        source_id="2",
        source_part="ppt/slides/slide1.xml",
    )


def test_group_identity_does_not_fall_through_to_child_metadata() -> None:
    group_element = ElementTree.fromstring(
        f'<p:grpSp xmlns:p="{_P}" xmlns:a="{_A}" xmlns:dx="urn:domoxml:canvas-ir:1">'
        "<p:nvGrpSpPr><p:cNvGrpSpPr/><p:nvPr/></p:nvGrpSpPr><p:grpSpPr/>"
        '<p:sp><p:nvSpPr><p:cNvPr id="7" name="child"/><p:cNvSpPr/><p:nvPr><p:extLst>'
        '<p:ext uri="{A6E4A7B1-9D9C-4E8F-A94B-22CB18A8D72F}">'
        '<dx:node id="child-id" sourceFormat="html" sourceId="child-source"/>'
        "</p:ext></p:extLst></p:nvPr></p:nvSpPr></p:sp></p:grpSp>"
    )
    group = GroupNode(
        box=Box(x=0, y=0, width=100, height=100),
        child_box=Box(x=0, y=0, width=100, height=100),
    )

    recovered = _with_pptx_identity(group, group_element, "ppt/slides/slide1.xml")

    assert recovered.node_id == "pptx-unknown"
    assert recovered.provenance == SourceProvenance(
        source_format="pptx",
        source_id="unknown",
        source_part="ppt/slides/slide1.xml",
    )


def test_reverse_coverage_records_flattened_group_components_and_source_loss() -> None:
    group_element = ElementTree.fromstring(
        f'<p:grpSp xmlns:p="{_P}"><p:nvGrpSpPr>'
        '<p:cNvPr id="9" name="group"/></p:nvGrpSpPr></p:grpSp>'
    )
    children = tuple(
        ShapeNode(box=Box(x=index * 100, y=0, width=100, height=100)) for index in range(2)
    )
    group = GroupNode(
        box=Box(x=0, y=0, width=200, height=100),
        child_box=Box(x=0, y=0, width=200, height=100),
        children=children,
    )

    coverage = _group_reverse_coverage(
        "ppt/slides/slide1.xml", group_element, group, has_preserved_children=False
    )

    assert coverage.representation is Representation.DECOMPOSED
    assert coverage.editability is Editability.COMPONENTS
    assert coverage.source_retention is SourceRetention.LOST
    assert coverage.output_count == 2

    with_preserved_child = _group_reverse_coverage(
        "ppt/slides/slide1.xml", group_element, group, has_preserved_children=True
    )
    assert with_preserved_child.source_retention is SourceRetention.LOST


def test_reverse_coverage_reports_connector_geometry_approximation() -> None:
    element = ElementTree.fromstring(
        f'<p:cxnSp xmlns:p="{_P}"><p:nvCxnSpPr>'
        '<p:cNvPr id="12" name="connector"/></p:nvCxnSpPr></p:cxnSp>'
    )

    coverage = _connector_reverse_coverage("ppt/slides/slide1.xml", element)

    assert coverage.representation is Representation.APPROXIMATED
    assert coverage.editability is Editability.SEMANTIC
    assert coverage.source_retention is SourceRetention.LOST


def test_resolves_theme_system_and_preset_colors() -> None:
    scheme = ElementTree.fromstring(
        f'<a:solidFill xmlns:a="{_A}"><a:schemeClr val="accent1">'
        '<a:alpha val="50000"/></a:schemeClr></a:solidFill>'
    )
    system = ElementTree.fromstring(
        f'<a:solidFill xmlns:a="{_A}"><a:sysClr val="window" lastClr="102030"/></a:solidFill>'
    )
    preset = ElementTree.fromstring(
        f'<a:solidFill xmlns:a="{_A}"><a:prstClr val="red"/></a:solidFill>'
    )
    assert rgba(scheme, {"accent1": "112233"}) == Rgba(r=17, g=34, b=51, a=0.5)
    assert rgba(system, {}) == Rgba(r=16, g=32, b=48)
    assert rgba(preset, {}) == Rgba(r=255, g=0, b=0)


def test_resolves_slide_theme_through_layout_master_chain_and_color_maps() -> None:
    package = OpcPackage.from_bytes(
        write_package(
            {
                "ppt/slides/slide1.xml": (
                    f'<p:sld xmlns:p="{_P}" xmlns:a="{_A}"><p:clrMapOvr>'
                    '<a:overrideClrMapping accent1="accent2"/></p:clrMapOvr></p:sld>'
                ),
                "ppt/slides/_rels/slide1.xml.rels": (
                    f'<Relationships xmlns="{_PKG_REL}"><Relationship Id="rId1" '
                    f'Type="{_R}/slideLayout" Target="../slideLayouts/slideLayout1.xml"/>'
                    "</Relationships>"
                ),
                "ppt/slideLayouts/slideLayout1.xml": (
                    f'<p:sldLayout xmlns:p="{_P}" xmlns:a="{_A}"><p:clrMapOvr>'
                    '<a:overrideClrMapping tx1="accent1"/></p:clrMapOvr></p:sldLayout>'
                ),
                "ppt/slideLayouts/_rels/slideLayout1.xml.rels": (
                    f'<Relationships xmlns="{_PKG_REL}"><Relationship Id="rId1" '
                    f'Type="{_R}/slideMaster" Target="../slideMasters/slideMaster1.xml"/>'
                    "</Relationships>"
                ),
                "ppt/slideMasters/slideMaster1.xml": (
                    f'<p:sldMaster xmlns:p="{_P}"><p:clrMap tx1="dk1" accent1="accent1"/>'
                    "</p:sldMaster>"
                ),
                "ppt/slideMasters/_rels/slideMaster1.xml.rels": (
                    f'<Relationships xmlns="{_PKG_REL}"><Relationship Id="rId1" '
                    f'Type="{_R}/theme" Target="../theme/theme1.xml"/></Relationships>'
                ),
                "ppt/theme/theme0.xml": (
                    f'<a:theme xmlns:a="{_A}"><a:themeElements><a:clrScheme name="wrong">'
                    '<a:accent1><a:srgbClr val="FFFFFF"/></a:accent1>'
                    "</a:clrScheme></a:themeElements></a:theme>"
                ),
                "ppt/theme/theme1.xml": (
                    f'<a:theme xmlns:a="{_A}"><a:themeElements><a:clrScheme name="right">'
                    '<a:dk1><a:srgbClr val="000000"/></a:dk1>'
                    '<a:accent1><a:srgbClr val="112233"/></a:accent1>'
                    '<a:accent2><a:srgbClr val="445566"/></a:accent2>'
                    "</a:clrScheme></a:themeElements></a:theme>"
                ),
            }
        )
    )
    colors = _slide_colors(package, "ppt/slides/slide1.xml")
    assert colors["tx1"] == "112233"
    assert colors["accent1"] == "445566"


def test_preserves_unsupported_slide_nodes_with_warning() -> None:
    package = OpcPackage.from_bytes(build_pptx([_sample_ir()], faces=[]))
    parts: dict[str, bytes | str] = {part: package.read(part) for part in package.parts}
    slide_part = "ppt/slides/slide1.xml"
    slide_xml = parts[slide_part]
    assert isinstance(slide_xml, bytes)
    parts[slide_part] = slide_xml.replace(
        b"</p:spTree>", b"<p:graphicFrame><p:nvGraphicFramePr/></p:graphicFrame></p:spTree>"
    )

    html = pptx_to_html(write_package(parts))

    assert len(html.preserved) == 1
    assert html.preserved[0].part == slide_part
    assert html.preserved[0].kind == "graphicFrame"
    assert "graphicFrame" in html.preserved[0].xml
    assert len(html.warnings) == 1
    assert "graphicFrame" in html.warnings[0].message


def test_attaches_element_crop_to_positioned_unsupported_node() -> None:
    package = OpcPackage.from_bytes(build_pptx([_sample_ir()], faces=[]))
    parts: dict[str, bytes | str] = {part: package.read(part) for part in package.parts}
    slide_part = "ppt/slides/slide1.xml"
    slide_xml = parts[slide_part]
    assert isinstance(slide_xml, bytes)
    frame = (
        f'<p:graphicFrame xmlns:p="{_P}" xmlns:a="{_A}">'
        '<p:nvGraphicFramePr><p:cNvPr id="20" name="unsupported"/>'
        "<p:cNvGraphicFramePr/><p:nvPr/></p:nvGraphicFramePr>"
        '<p:xfrm><a:off x="3048000" y="952500"/>'
        '<a:ext cx="1905000" cy="1428750"/></p:xfrm>'
        '<a:graphic><a:graphicData uri="urn:example:unsupported"/></a:graphic>'
        "</p:graphicFrame>"
    ).encode()
    parts[slide_part] = slide_xml.replace(b"</p:spTree>", frame + b"</p:spTree>")
    rendered = BytesIO()
    Image.new("RGB", (1280, 720), "#2A7F62").save(rendered, "PNG")

    result = read_pptx_result(write_package(parts), fallback_pngs=(rendered.getvalue(),))

    [preserved] = [node for node in result.slides[0].contents if isinstance(node, PreservedNode)]
    assert preserved.node_id == "pptx-20"
    assert preserved.fallback is not None
    with Image.open(BytesIO(preserved.fallback.data)) as crop:
        assert crop.size == (200, 150)
        pixel = crop.convert("RGB").getpixel((100, 75))
        assert isinstance(pixel, tuple)
        assert pixel == (42, 127, 98)
    [fragment] = [fragment for fragment in result.preserved if fragment.kind == "graphicFrame"]
    assert fragment.owner_node_id == preserved.node_id
    [coverage] = [
        item
        for item in result.coverage.items
        if item.element == "ppt/slides/slide1.xml:graphicFrame#20"
    ]
    assert coverage.representation is Representation.ELEMENT_LAYER
    assert coverage.editability is Editability.LAYERS
    assert coverage.source_retention is SourceRetention.ATTACHED
    assert coverage.output_count == 1
    assert coverage.raster_area_emu2 == 1_905_000 * 1_428_750


def test_attached_source_without_renderer_is_reported_as_visually_failed() -> None:
    package = OpcPackage.from_bytes(build_pptx([_sample_ir()], faces=[]))
    parts: dict[str, bytes | str] = {part: package.read(part) for part in package.parts}
    slide_part = "ppt/slides/slide1.xml"
    slide_xml = parts[slide_part]
    assert isinstance(slide_xml, bytes)
    frame = (
        f'<p:graphicFrame xmlns:p="{_P}" xmlns:a="{_A}">'
        '<p:nvGraphicFramePr><p:cNvPr id="20" name="unsupported"/>'
        "<p:cNvGraphicFramePr/><p:nvPr/></p:nvGraphicFramePr>"
        '<p:xfrm><a:off x="3048000" y="952500"/>'
        '<a:ext cx="1905000" cy="1428750"/></p:xfrm>'
        '<a:graphic><a:graphicData uri="urn:example:unsupported"/></a:graphic>'
        "</p:graphicFrame>"
    ).encode()
    parts[slide_part] = slide_xml.replace(b"</p:spTree>", frame + b"</p:spTree>")

    result = read_pptx_result(write_package(parts))

    [coverage] = [
        item
        for item in result.coverage.items
        if item.element == "ppt/slides/slide1.xml:graphicFrame#20"
    ]
    assert coverage.representation is Representation.FAILED
    assert coverage.editability is Editability.NONE
    assert coverage.source_retention is SourceRetention.ATTACHED
    assert coverage.output_count == 0
    assert coverage.raster_area_emu2 == 0


def test_rejects_mismatched_reverse_fallback_page_count() -> None:
    pptx = build_pptx([_sample_ir()], faces=[])

    try:
        read_pptx_result(pptx, fallback_pngs=())
    except ValueError as error:
        assert str(error) == "fallback PNG count 0 does not match slide count 1"
    else:  # pragma: no cover - assertion helper without a pytest dependency
        raise AssertionError("expected mismatched fallback count to fail")


def test_malformed_preserved_graph_still_gets_a_visual_layer() -> None:
    package = OpcPackage.from_bytes(build_pptx([_sample_ir()], faces=[]))
    parts: dict[str, bytes | str] = {part: package.read(part) for part in package.parts}
    slide_part = "ppt/slides/slide1.xml"
    slide_xml = parts[slide_part]
    assert isinstance(slide_xml, bytes)
    frame = (
        f'<p:graphicFrame xmlns:p="{_P}" xmlns:a="{_A}" xmlns:r="{_R}">'
        '<p:nvGraphicFramePr><p:cNvPr id="20" name="malformed"/>'
        "<p:cNvGraphicFramePr/><p:nvPr/></p:nvGraphicFramePr>"
        '<p:xfrm><a:off x="3048000" y="952500"/>'
        '<a:ext cx="1905000" cy="1428750"/></p:xfrm>'
        '<a:graphic><a:graphicData uri="urn:example:unsupported">'
        '<a:ext r:id="rIdMissing"/></a:graphicData></a:graphic>'
        "</p:graphicFrame>"
    ).encode()
    parts[slide_part] = slide_xml.replace(b"</p:spTree>", frame + b"</p:spTree>")
    rendered = BytesIO()
    Image.new("RGB", (1280, 720), "#2A7F62").save(rendered, "PNG")

    result = read_pptx_result(write_package(parts), fallback_pngs=(rendered.getvalue(),))

    [layer] = [
        node
        for node in result.slides[0].contents
        if isinstance(node, ShapeNode)
        and isinstance(node.fill, PictureFill)
        and node.fill.raster_role == "pptx-source-fallback"
    ]
    assert layer.node_id == "pptx-20"
    [fragment] = [fragment for fragment in result.preserved if fragment.kind == "graphicFrame"]
    assert fragment.owner_node_id == layer.node_id
    assert "visual retained as an element layer" in result.warnings[-1].message
    [coverage] = [
        item
        for item in result.coverage.items
        if item.element == "ppt/slides/slide1.xml:graphicFrame#20"
    ]
    assert coverage.representation is Representation.ELEMENT_LAYER
    assert coverage.source_retention is SourceRetention.DETACHED
    assert coverage.raster_area_emu2 == 1_905_000 * 1_428_750


def test_reads_native_powerpoint_picture() -> None:
    image = BytesIO()
    Image.new("RGB", (10, 10), "#4472C4").save(image, format="PNG")
    image.seek(0)
    deck = PptxPresentation()
    slide = deck.slides.add_slide(deck.slide_layouts[6])
    slide.shapes.add_picture(image, Inches(1), Inches(1), Inches(2), Inches(1))
    pptx = BytesIO()
    deck.save(pptx)

    [result] = read_pptx(pptx.getvalue())

    assert len(result.shapes) == 1
    fill = result.shapes[0].fill
    assert isinstance(fill, PictureFill)
    assert fill.ext == "png"


def _two_slide_deck(run_xml: str, extra_slide1_rels: str = "") -> bytes:
    """A minimal two-slide deck where slide1's only shape carries ``run_xml`` in its text body.

    ``extra_slide1_rels`` injects hyperlink relationships referenced by the run XML."""
    sp = (
        f'<p:sp xmlns:p="{_P}" xmlns:a="{_A}" xmlns:r="{_R}"><p:nvSpPr>'
        '<p:cNvPr id="2" name="s"/><p:cNvSpPr/><p:nvPr/></p:nvSpPr>'
        '<p:spPr><a:xfrm><a:off x="0" y="0"/><a:ext cx="3000000" cy="1000000"/></a:xfrm>'
        '<a:prstGeom prst="rect"><a:avLst/></a:prstGeom></p:spPr>'
        f'<p:txBody><a:bodyPr/><a:lstStyle/><a:p><a:pPr algn="l"/>{run_xml}</a:p></p:txBody></p:sp>'
    )
    parts: dict[str, bytes | str] = {
        "[Content_Types].xml": (
            '<?xml version="1.0"?><Types '
            'xmlns="http://schemas.openxmlformats.org/package/2006/content-types">'
            '<Default Extension="rels" '
            'ContentType="application/vnd.openxmlformats-package.relationships+xml"/>'
            '<Default Extension="xml" ContentType="application/xml"/></Types>'
        ),
        "_rels/.rels": (
            f'<Relationships xmlns="{_PKG_REL}"><Relationship Id="rId1" '
            f'Type="{_R}/officeDocument" Target="ppt/presentation.xml"/></Relationships>'
        ),
        "ppt/presentation.xml": (
            f'<p:presentation xmlns:p="{_P}" xmlns:r="{_R}"><p:sldIdLst>'
            '<p:sldId id="256" r:id="rId1"/><p:sldId id="257" r:id="rId2"/></p:sldIdLst>'
            '<p:sldSz cx="12192000" cy="6858000"/></p:presentation>'
        ),
        "ppt/_rels/presentation.xml.rels": (
            f'<Relationships xmlns="{_PKG_REL}">'
            f'<Relationship Id="rId1" Type="{_R}/slide" Target="slides/slide1.xml"/>'
            f'<Relationship Id="rId2" Type="{_R}/slide" Target="slides/slide2.xml"/>'
            "</Relationships>"
        ),
        "ppt/slides/slide1.xml": (
            f'<p:sld xmlns:p="{_P}" xmlns:a="{_A}"><p:cSld><p:spTree>'
            '<p:nvGrpSpPr><p:cNvPr id="1" name=""/><p:cNvGrpSpPr/><p:nvPr/></p:nvGrpSpPr>'
            f"<p:grpSpPr/>{sp}</p:spTree></p:cSld></p:sld>"
        ),
        "ppt/slides/_rels/slide1.xml.rels": (
            f'<Relationships xmlns="{_PKG_REL}">{extra_slide1_rels}</Relationships>'
        ),
        "ppt/slides/slide2.xml": (
            f'<p:sld xmlns:p="{_P}" xmlns:a="{_A}"><p:cSld><p:spTree>'
            '<p:nvGrpSpPr><p:cNvPr id="1" name=""/><p:cNvGrpSpPr/><p:nvPr/></p:nvGrpSpPr>'
            "<p:grpSpPr/></p:spTree></p:cSld></p:sld>"
        ),
        "ppt/slides/_rels/slide2.xml.rels": f'<Relationships xmlns="{_PKG_REL}"/>',
    }
    return write_package(parts)


def test_reverse_reads_underline_strike_caps_and_letter_spacing() -> None:
    run = (
        '<a:r><a:rPr lang="en-US" sz="1800" u="sng" strike="sngStrike" cap="all" spc="200">'
        '<a:latin typeface="Inter"/></a:rPr><a:t>x</a:t></a:r>'
    )
    [slide, _] = read_pptx(_two_slide_deck(run))
    text = slide.shapes[0].text
    assert text is not None
    r = text.paragraphs[0].runs[0]
    assert r.underline is True
    assert r.strike is True
    assert r.caps == "all"
    assert r.letter_spacing_pt == 2.0


def test_reverse_reads_small_caps() -> None:
    run = (
        '<a:r><a:rPr lang="en-US" sz="1800" cap="small"><a:latin typeface="Inter"/></a:rPr>'
        "<a:t>x</a:t></a:r>"
    )
    [slide, _] = read_pptx(_two_slide_deck(run))
    assert slide.shapes[0].text is not None
    assert slide.shapes[0].text.paragraphs[0].runs[0].caps == "small"


def test_reverse_reads_external_hyperlink() -> None:
    run = (
        f'<a:r xmlns:r="{_R}"><a:rPr lang="en-US" sz="1800"><a:hlinkClick r:id="rIdL"/>'
        '<a:latin typeface="Inter"/></a:rPr><a:t>x</a:t></a:r>'
    )
    rels = (
        f'<Relationship Id="rIdL" Type="{_R}/hyperlink" '
        'Target="https://example.com" TargetMode="External"/>'
    )
    [slide, _] = read_pptx(_two_slide_deck(run, rels))
    assert slide.shapes[0].text is not None
    link = slide.shapes[0].text.paragraphs[0].runs[0].hyperlink
    assert link is not None and link.url == "https://example.com"


def test_reverse_reads_slide_jump_hyperlink_into_index() -> None:
    run = (
        f'<a:r xmlns:r="{_R}"><a:rPr lang="en-US" sz="1800">'
        '<a:hlinkClick r:id="rIdJ" action="ppaction://hlinksldjump"/>'
        '<a:latin typeface="Inter"/></a:rPr><a:t>x</a:t></a:r>'
    )
    rels = f'<Relationship Id="rIdJ" Type="{_R}/slide" Target="slide2.xml"/>'
    [slide, _] = read_pptx(_two_slide_deck(run, rels))
    assert slide.shapes[0].text is not None
    link = slide.shapes[0].text.paragraphs[0].runs[0].hyperlink
    assert link is not None and link.slide_index == 1  # slide2.xml is the 2nd slide (index 1)


# --------------------------------------------------------------------------- paragraph pPr reverse


def _para_deck(para_xml: str) -> bytes:
    """A minimal one-slide deck whose only shape has one paragraph ``para_xml``."""
    sp = (
        f'<p:sp xmlns:p="{_P}" xmlns:a="{_A}" xmlns:r="{_R}"><p:nvSpPr>'
        '<p:cNvPr id="2" name="s"/><p:cNvSpPr/><p:nvPr/></p:nvSpPr>'
        '<p:spPr><a:xfrm><a:off x="0" y="0"/><a:ext cx="3000000" cy="1000000"/></a:xfrm>'
        '<a:prstGeom prst="rect"><a:avLst/></a:prstGeom></p:spPr>'
        f"<p:txBody><a:bodyPr/><a:lstStyle/>{para_xml}</p:txBody></p:sp>"
    )
    parts: dict[str, bytes | str] = {
        "[Content_Types].xml": (
            '<?xml version="1.0"?><Types '
            'xmlns="http://schemas.openxmlformats.org/package/2006/content-types">'
            '<Default Extension="rels" '
            'ContentType="application/vnd.openxmlformats-package.relationships+xml"/>'
            '<Default Extension="xml" ContentType="application/xml"/></Types>'
        ),
        "_rels/.rels": (
            f'<Relationships xmlns="{_PKG_REL}"><Relationship Id="rId1" '
            f'Type="{_R}/officeDocument" Target="ppt/presentation.xml"/></Relationships>'
        ),
        "ppt/presentation.xml": (
            f'<p:presentation xmlns:p="{_P}" xmlns:r="{_R}"><p:sldIdLst>'
            '<p:sldId id="256" r:id="rId1"/></p:sldIdLst>'
            '<p:sldSz cx="12192000" cy="6858000"/></p:presentation>'
        ),
        "ppt/_rels/presentation.xml.rels": (
            f'<Relationships xmlns="{_PKG_REL}">'
            f'<Relationship Id="rId1" Type="{_R}/slide" Target="slides/slide1.xml"/>'
            "</Relationships>"
        ),
        "ppt/slides/slide1.xml": (
            f'<p:sld xmlns:p="{_P}" xmlns:a="{_A}"><p:cSld><p:spTree>'
            '<p:nvGrpSpPr><p:cNvPr id="1" name=""/><p:cNvGrpSpPr/><p:nvPr/></p:nvGrpSpPr>'
            f"<p:grpSpPr/>{sp}</p:spTree></p:cSld></p:sld>"
        ),
        "ppt/slides/_rels/slide1.xml.rels": f'<Relationships xmlns="{_PKG_REL}"/>',
    }
    return write_package(parts)


def _first_para(pptx: bytes) -> TextParagraph:
    [slide] = read_pptx(pptx)
    assert slide.shapes[0].text is not None
    return slide.shapes[0].text.paragraphs[0]


def test_reverse_reads_lnspc_pct() -> None:
    """<a:lnSpc><a:spcPct val="150000"/> → LineSpacing(percent=1.5)."""
    para_xml = (
        f'<a:p xmlns:a="{_A}"><a:pPr algn="l">'
        '<a:lnSpc><a:spcPct val="150000"/></a:lnSpc>'
        '</a:pPr><a:r><a:rPr lang="en-US" sz="1200"/><a:t>x</a:t></a:r></a:p>'
    )
    para = _first_para(_para_deck(para_xml))
    assert para.line_spacing is not None
    assert para.line_spacing.percent == 1.5


def test_reverse_reads_lnspc_pts() -> None:
    """<a:lnSpc><a:spcPts val="1800"/> → LineSpacing(points=18.0)."""
    para_xml = (
        f'<a:p xmlns:a="{_A}"><a:pPr algn="l">'
        '<a:lnSpc><a:spcPts val="1800"/></a:lnSpc>'
        '</a:pPr><a:r><a:rPr lang="en-US" sz="1200"/><a:t>x</a:t></a:r></a:p>'
    )
    para = _first_para(_para_deck(para_xml))
    assert para.line_spacing is not None
    assert para.line_spacing.points == 18.0


def test_reverse_reads_spc_bef_aft() -> None:
    """<a:spcBef><a:spcPts val="900"/> → space_before_pt=9.0, etc."""
    para_xml = (
        f'<a:p xmlns:a="{_A}"><a:pPr algn="l">'
        '<a:spcBef><a:spcPts val="900"/></a:spcBef>'
        '<a:spcAft><a:spcPts val="1800"/></a:spcAft>'
        '</a:pPr><a:r><a:rPr lang="en-US" sz="1200"/><a:t>x</a:t></a:r></a:p>'
    )
    para = _first_para(_para_deck(para_xml))
    assert para.space_before_pt == 9.0
    assert para.space_after_pt == 18.0


def test_reverse_reads_mar_l_and_indent() -> None:
    """marL="457200" → left_margin_pt=36.0; indent="228600" → indent_pt=18.0."""
    para_xml = (
        f'<a:p xmlns:a="{_A}"><a:pPr algn="l" marL="457200" indent="228600"/>'
        '<a:r><a:rPr lang="en-US" sz="1200"/><a:t>x</a:t></a:r></a:p>'
    )
    para = _first_para(_para_deck(para_xml))
    assert abs(para.left_margin_pt - 36.0) < 0.1  # 457200 / 12700 = 36
    assert abs(para.indent_pt - 18.0) < 0.1  # 228600 / 12700 = 18


def test_reverse_reads_bu_char() -> None:
    """<a:buChar char="•"/> → CharBullet(char="•")."""
    para_xml = (
        f'<a:p xmlns:a="{_A}"><a:pPr algn="l">'
        '<a:buChar char="&#x2022;"/>'
        '</a:pPr><a:r><a:rPr lang="en-US" sz="1200"/><a:t>item</a:t></a:r></a:p>'
    )
    para = _first_para(_para_deck(para_xml))
    assert isinstance(para.bullet, CharBullet)
    assert para.bullet.char == "•"


def test_reverse_reads_bu_autonum() -> None:
    """<a:buAutoNum type="arabicPeriod" startAt="1"/> → AutoNumberBullet(...)."""
    para_xml = (
        f'<a:p xmlns:a="{_A}"><a:pPr algn="l">'
        '<a:buAutoNum type="arabicPeriod" startAt="1"/>'
        '</a:pPr><a:r><a:rPr lang="en-US" sz="1200"/><a:t>item</a:t></a:r></a:p>'
    )
    para = _first_para(_para_deck(para_xml))
    assert isinstance(para.bullet, AutoNumberBullet)
    assert para.bullet.scheme == "arabicPeriod"
    assert para.bullet.start_at == 1


def test_reverse_reads_nested_level() -> None:
    """lvl="2" → level=2."""
    para_xml = (
        f'<a:p xmlns:a="{_A}"><a:pPr algn="l" lvl="2">'
        '<a:buChar char="&#x25AA;"/>'
        '</a:pPr><a:r><a:rPr lang="en-US" sz="1200"/><a:t>deep</a:t></a:r></a:p>'
    )
    para = _first_para(_para_deck(para_xml))
    assert para.level == 2
