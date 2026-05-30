"""End-to-end: HTML → editable .pptx. Requires Chromium: ``playwright install chromium``."""

from __future__ import annotations

import io
from pathlib import Path

import pytest
from pptx import Presentation as PptxRead  # test-only validator

from domoxml import Presentation, Slide
from domoxml.types import OutputFormat, SlideSize

pytestmark = pytest.mark.integration


def test_render_pptx_is_editable_and_keeps_the_text() -> None:
    deck = Presentation(size=SlideSize.WIDE_16_9)
    deck.add(
        Slide(
            html="<div style='background:rgb(79,70,229);color:#fff;font-size:40px'>Driftwood</div>"
        )
    )
    result = deck.render({OutputFormat.PPTX})

    assert result.pptx is not None and result.pptx[:2] == b"PK"
    prs = PptxRead(io.BytesIO(result.pptx))
    assert len(prs.slides) == 1
    texts: list[str] = []
    for shape in prs.slides[0].shapes:
        if shape.has_text_frame:
            texts.append(shape.text_frame.text)  # pyright: ignore  (python-pptx stubs)
    assert any("Driftwood" in text for text in texts)  # native text, not a flat image


def test_save_pptx_writes_an_openable_file(tmp_path: Path) -> None:
    deck = Presentation()
    deck.add(Slide(html="<p>hello</p>"))
    out = tmp_path / "deck.pptx"
    deck.save_pptx(out)
    assert out.exists() and out.stat().st_size > 0
    PptxRead(str(out))  # opens without raising
