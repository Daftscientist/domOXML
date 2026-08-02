"""Slides-backend unit tests: build a .pptx from IR and prove it's valid (no browser)."""

from __future__ import annotations

import io
import zipfile
from xml.etree import ElementTree

import pytest
from PIL import Image
from pptx import Presentation as PptxRead  # test-only validator

from domoxml.core.ir.model import (
    AutoNumberBullet,
    Box,
    CharBullet,
    Connector,
    GroupNode,
    Hyperlink,
    Line,
    LineSpacing,
    PictureFill,
    Point,
    PortableFallback,
    PreservationPart,
    PreservationPayload,
    PreservedNode,
    Rgba,
    ShapeNode,
    SlideIR,
    SolidFill,
    SrcRect,
    TextBody,
    TextParagraph,
    TextRun,
    Transform,
)
from domoxml.core.opc import OpcPackage, write_package
from domoxml.slides import build_pptx, read_pptx, read_pptx_result
from domoxml.types import Editability, Representation, SourceRetention

_A = "http://schemas.openxmlformats.org/drawingml/2006/main"


def _slide_xml(pptx: bytes, name: str = "ppt/slides/slide1.xml") -> str:
    with zipfile.ZipFile(io.BytesIO(pptx)) as archive:
        return archive.read(name).decode("utf-8")


def _slide_rels(pptx: bytes, name: str = "ppt/slides/_rels/slide1.xml.rels") -> str:
    with zipfile.ZipFile(io.BytesIO(pptx)) as archive:
        return archive.read(name).decode("utf-8")


def _decorated_run(**kw: object) -> SlideIR:
    return SlideIR(
        width=12_192_000,
        height=6_858_000,
        shapes=(
            ShapeNode(
                box=Box(x=0, y=0, width=3_000_000, height=1_000_000),
                text=TextBody(
                    paragraphs=(
                        TextParagraph(
                            runs=(TextRun(text="x", font_family="Inter", size_pt=18, **kw),)  # type: ignore[arg-type]
                        ),
                    )
                ),
            ),
        ),
    )


def _sample_ir() -> SlideIR:
    return SlideIR(
        width=12_192_000,
        height=6_858_000,
        shapes=(
            ShapeNode(
                box=Box(x=914_400, y=914_400, width=3_657_600, height=1_828_800),
                fill=SolidFill(color=Rgba(r=79, g=70, b=229)),
                corner_radius_emu=76_200,
                text=TextBody(
                    paragraphs=(
                        TextParagraph(
                            runs=(
                                TextRun(
                                    text="Driftwood",
                                    font_family="Inter",
                                    size_pt=24.0,
                                    bold=True,
                                    color=Rgba(r=255, g=255, b=255),
                                ),
                            ),
                            align="center",
                        ),
                    )
                ),
            ),
            ShapeNode(box=Box(x=0, y=0, width=100, height=100)),  # plain, no fill/text
        ),
    )


def test_build_pptx_is_a_zip_with_required_parts() -> None:
    data = build_pptx([_sample_ir()])
    assert data[:2] == b"PK"
    with zipfile.ZipFile(io.BytesIO(data)) as archive:
        names = set(archive.namelist())
    required = {
        "[Content_Types].xml",
        "_rels/.rels",
        "ppt/presentation.xml",
        "ppt/slideMasters/slideMaster1.xml",
        "ppt/slideLayouts/slideLayout1.xml",
        "ppt/theme/theme1.xml",
        "ppt/slides/slide1.xml",
    }
    assert required <= names


def test_build_pptx_opens_and_keeps_text_editable() -> None:
    prs = PptxRead(io.BytesIO(build_pptx([_sample_ir()])))
    assert len(prs.slides) == 1
    texts: list[str] = []
    for shape in prs.slides[0].shapes:
        if shape.has_text_frame:
            texts.append(shape.text_frame.text)  # pyright: ignore  (python-pptx stubs)
    assert "Driftwood" in texts  # real editable text run, not a rasterised image


def test_native_group_round_trips_through_pptx_as_one_editable_group() -> None:
    children = (
        ShapeNode(
            box=Box(x=100_000, y=200_000, width=800_000, height=500_000),
            fill=SolidFill(color=Rgba(r=239, g=68, b=68)),
        ),
        ShapeNode(
            box=Box(x=1_000_000, y=350_000, width=600_000, height=700_000),
            fill=SolidFill(color=Rgba(r=37, g=99, b=235)),
        ),
    )
    group = GroupNode(
        node_id="group-1",
        box=Box(x=2_000_000, y=1_000_000, width=3_200_000, height=2_100_000),
        child_box=Box(x=100_000, y=200_000, width=1_500_000, height=850_000),
        children=children,
    )

    [recovered] = read_pptx(
        build_pptx(
            [SlideIR(width=12_192_000, height=6_858_000, contents=(group,))],
            faces=[],
        )
    )[0].contents

    assert isinstance(recovered, GroupNode)
    assert recovered.node_id == "group-1"
    assert recovered.box == group.box
    assert recovered.child_box == group.child_box
    recovered_shapes = [child for child in recovered.children if isinstance(child, ShapeNode)]
    expected_shapes = [child for child in group.children if isinstance(child, ShapeNode)]
    assert len(recovered_shapes) == len(recovered.children)
    assert [child.box for child in recovered_shapes] == [child.box for child in expected_shapes]
    assert [child.fill for child in recovered_shapes] == [child.fill for child in expected_shapes]


def test_native_group_reverse_coverage_reports_retained_group_semantics() -> None:
    group = GroupNode(
        box=Box(x=2_000_000, y=1_000_000, width=3_200_000, height=2_100_000),
        child_box=Box(x=0, y=0, width=1_600_000, height=1_050_000),
        children=(
            ShapeNode(
                box=Box(x=0, y=0, width=800_000, height=500_000),
                fill=SolidFill(color=Rgba(r=239, g=68, b=68)),
            ),
            ShapeNode(
                box=Box(x=900_000, y=350_000, width=600_000, height=700_000),
                fill=SolidFill(color=Rgba(r=37, g=99, b=235)),
            ),
        ),
    )

    result = read_pptx_result(
        build_pptx(
            [SlideIR(width=12_192_000, height=6_858_000, contents=(group,))],
            faces=[],
        )
    )

    [coverage] = result.coverage.items
    assert coverage.representation is Representation.NATIVE
    assert coverage.editability is Editability.SEMANTIC
    assert coverage.source_retention is SourceRetention.NOT_REQUIRED
    assert coverage.output_count == 1


def test_group_with_unsupported_child_retains_source_with_honest_rasterized_visual() -> None:
    group = GroupNode(
        box=Box(x=2_000_000, y=1_000_000, width=3_200_000, height=2_100_000),
        child_box=Box(x=0, y=0, width=1_600_000, height=1_050_000),
        children=(
            ShapeNode(
                box=Box(x=0, y=0, width=800_000, height=500_000),
                fill=SolidFill(color=Rgba(r=239, g=68, b=68)),
            ),
        ),
    )
    package = OpcPackage.from_bytes(
        build_pptx(
            [
                SlideIR(
                    width=12_192_000,
                    height=6_858_000,
                    contents=(
                        group,
                        ShapeNode(
                            box=Box(
                                x=3_000_000,
                                y=1_500_000,
                                width=2_000_000,
                                height=1_500_000,
                            ),
                            fill=SolidFill(color=Rgba(r=245, g=158, b=11, a=0.6)),
                        ),
                    ),
                )
            ],
            faces=[],
        )
    )
    parts: dict[str, bytes | str] = {part: package.read(part) for part in package.parts}
    slide_part = "ppt/slides/slide1.xml"
    root = ElementTree.fromstring(parts[slide_part])
    native_group = root.find(".//{http://schemas.openxmlformats.org/presentationml/2006/main}grpSp")
    assert native_group is not None
    ElementTree.SubElement(
        native_group,
        "{http://schemas.openxmlformats.org/presentationml/2006/main}contentPart",
    )
    parts[slide_part] = ElementTree.tostring(root)
    source = write_package(parts)
    fallback = io.BytesIO()
    Image.new("RGB", (1280, 720), "#f8fafc").save(fallback, "PNG")

    result = read_pptx_result(source, fallback_pngs=[fallback.getvalue()])

    preserved, sibling = result.slides[0].contents
    assert isinstance(preserved, PreservedNode)
    assert isinstance(sibling, ShapeNode)
    assert preserved.payload.kind == "grpSp"
    assert preserved.fallback is None
    assert result.slides[0].renderer_fallback is not None
    assert result.slides[0].renderer_fallback_owner_node_id == preserved.node_id
    rasterized = next(
        item for item in result.coverage.items if item.representation is Representation.RASTERIZED
    )
    assert rasterized.editability is Editability.NONE
    assert rasterized.source_retention is SourceRetention.ATTACHED
    assert rasterized.raster_area_emu2 == 12_192_000 * 6_858_000
    rebuilt_package = OpcPackage.from_bytes(build_pptx(list(result.slides), faces=[]))
    rebuilt = rebuilt_package.read(slide_part)
    assert b"<p:grpSp" in rebuilt
    assert b"<p:contentPart" in rebuilt
    root = ElementTree.fromstring(rebuilt)
    ids = [
        element.get("id")
        for element in root.findall(
            ".//{http://schemas.openxmlformats.org/presentationml/2006/main}cNvPr"
        )
    ]
    assert len(ids) == len(set(ids))


def test_transformed_group_stays_owned_until_group_transform_html_is_exact() -> None:
    group = GroupNode(
        box=Box(x=2_000_000, y=1_000_000, width=3_200_000, height=2_100_000),
        child_box=Box(x=0, y=0, width=1_600_000, height=1_050_000),
        children=(
            ShapeNode(
                box=Box(x=0, y=0, width=800_000, height=500_000),
                fill=SolidFill(color=Rgba(r=239, g=68, b=68)),
            ),
        ),
        transform=Transform(rotation_deg=15),
    )
    source = build_pptx([SlideIR(width=12_192_000, height=6_858_000, contents=(group,))], faces=[])
    fallback = io.BytesIO()
    Image.new("RGB", (1280, 720), "#f8fafc").save(fallback, "PNG")

    result = read_pptx_result(source, fallback_pngs=[fallback.getvalue()])

    [preserved] = result.slides[0].contents
    assert isinstance(preserved, PreservedNode)
    assert preserved.payload.kind == "grpSp"
    assert preserved.fallback is None
    assert result.slides[0].renderer_fallback is not None
    assert result.slides[0].renderer_fallback_owner_node_id == preserved.node_id
    [coverage] = result.coverage.items
    assert coverage.representation is Representation.RASTERIZED
    assert coverage.editability is Editability.NONE
    assert coverage.source_retention is SourceRetention.ATTACHED

    without_renderer = read_pptx_result(source)
    [preserved_without_renderer] = without_renderer.slides[0].contents
    assert isinstance(preserved_without_renderer, PreservedNode)
    assert preserved_without_renderer.fallback is None
    assert without_renderer.slides[0].renderer_fallback is None
    [failed] = without_renderer.coverage.items
    assert failed.representation is Representation.FAILED
    assert failed.editability is Editability.NONE
    assert failed.source_retention is SourceRetention.ATTACHED


def test_nested_group_retains_source_until_nested_html_reconstruction_is_proven() -> None:
    inner = GroupNode(
        box=Box(x=0, y=0, width=1_600_000, height=1_050_000),
        child_box=Box(x=0, y=0, width=1_600_000, height=1_050_000),
        children=(
            ShapeNode(
                box=Box(x=0, y=0, width=800_000, height=500_000),
                fill=SolidFill(color=Rgba(r=239, g=68, b=68)),
            ),
        ),
    )
    outer = GroupNode(
        box=Box(x=2_000_000, y=1_000_000, width=3_200_000, height=2_100_000),
        child_box=Box(x=0, y=0, width=1_600_000, height=1_050_000),
        children=(inner,),
    )
    source = build_pptx([SlideIR(width=12_192_000, height=6_858_000, contents=(outer,))], faces=[])
    fallback = io.BytesIO()
    Image.new("RGB", (1280, 720), "#f8fafc").save(fallback, "PNG")

    result = read_pptx_result(source, fallback_pngs=[fallback.getvalue()])

    [preserved] = result.slides[0].contents
    assert isinstance(preserved, PreservedNode)
    assert preserved.payload.kind == "grpSp"
    assert preserved.fallback is None
    assert result.slides[0].renderer_fallback is not None
    assert result.slides[0].renderer_fallback_owner_node_id == preserved.node_id
    [coverage] = result.coverage.items
    assert coverage.representation is Representation.RASTERIZED
    assert coverage.source_retention is SourceRetention.ATTACHED


def test_multiple_unsupported_groups_keep_one_owner_and_explicit_remaining_debt() -> None:
    first = GroupNode(
        node_id="unsupported-group-1",
        box=Box(x=1_000_000, y=1_000_000, width=2_000_000, height=1_500_000),
        child_box=Box(x=0, y=0, width=1_000_000, height=750_000),
        children=(
            ShapeNode(
                box=Box(x=0, y=0, width=1_000_000, height=750_000),
                fill=SolidFill(color=Rgba(r=239, g=68, b=68)),
            ),
        ),
        transform=Transform(rotation_deg=10),
    )
    second = first.model_copy(
        update={
            "node_id": "unsupported-group-2",
            "box": Box(x=4_000_000, y=2_500_000, width=2_000_000, height=1_500_000),
            "transform": Transform(rotation_deg=-12),
        }
    )
    source = build_pptx(
        [SlideIR(width=12_192_000, height=6_858_000, contents=(first, second))],
        faces=[],
    )
    fallback = io.BytesIO()
    Image.new("RGB", (1280, 720), "#f8fafc").save(fallback, "PNG")

    result = read_pptx_result(source, fallback_pngs=[fallback.getvalue()])

    first_preserved, second_preserved = result.slides[0].contents
    assert isinstance(first_preserved, PreservedNode)
    assert isinstance(second_preserved, PreservedNode)
    assert first_preserved.node_id != second_preserved.node_id
    assert result.slides[0].renderer_fallback_owner_node_id == first_preserved.node_id
    representations = [item.representation for item in result.coverage.items]
    assert representations.count(Representation.RASTERIZED) == 1
    assert representations.count(Representation.FAILED) == 1
    failed = next(
        item for item in result.coverage.items if item.representation is Representation.FAILED
    )
    assert failed.source_retention is SourceRetention.ATTACHED


def test_preserved_group_rewrites_connector_shape_id_references() -> None:
    group = GroupNode(
        box=Box(x=1_000_000, y=1_000_000, width=3_000_000, height=2_000_000),
        child_box=Box(x=0, y=0, width=3_000_000, height=2_000_000),
        children=(
            ShapeNode(box=Box(x=0, y=0, width=500_000, height=500_000)),
            ShapeNode(box=Box(x=2_000_000, y=1_000_000, width=500_000, height=500_000)),
            Connector(
                start=Point(x=500_000, y=500_000),
                end=Point(x=2_000_000, y=1_000_000),
                line=Line(color=Rgba(r=15, g=23, b=42), width_emu=12_700),
            ),
        ),
    )
    package = OpcPackage.from_bytes(
        build_pptx(
            [
                SlideIR(
                    width=12_192_000,
                    height=6_858_000,
                    contents=(
                        group,
                        ShapeNode(box=Box(x=5_000_000, y=1_000_000, width=500_000, height=500_000)),
                    ),
                )
            ],
            faces=[],
        )
    )
    parts: dict[str, bytes | str] = {part: package.read(part) for part in package.parts}
    slide_part = "ppt/slides/slide1.xml"
    root = ElementTree.fromstring(parts[slide_part])
    namespace = "http://schemas.openxmlformats.org/presentationml/2006/main"
    drawing_namespace = "http://schemas.openxmlformats.org/drawingml/2006/main"
    native_group = root.find(f".//{{{namespace}}}grpSp")
    assert native_group is not None
    non_visuals = native_group.findall(f".//{{{namespace}}}cNvPr")
    for non_visual, shape_id in zip(non_visuals, (100, 220, 340, 460), strict=True):
        non_visual.set("id", str(shape_id))
    connection_properties = native_group.find(f".//{{{namespace}}}cNvCxnSpPr")
    assert connection_properties is not None
    ElementTree.SubElement(
        connection_properties,
        f"{{{drawing_namespace}}}stCxn",
        {"id": "220", "idx": "0"},
    )
    ElementTree.SubElement(
        connection_properties,
        f"{{{drawing_namespace}}}endCxn",
        {"id": "340", "idx": "0"},
    )
    parts[slide_part] = ElementTree.tostring(root)
    source = write_package(parts)
    fallback = io.BytesIO()
    Image.new("RGB", (1280, 720), "white").save(fallback, "PNG")

    result = read_pptx_result(source, fallback_pngs=[fallback.getvalue()])
    rebuilt = OpcPackage.from_bytes(build_pptx(list(result.slides), faces=[])).read(slide_part)
    rebuilt_root = ElementTree.fromstring(rebuilt)
    rebuilt_group = rebuilt_root.find(f".//{{{namespace}}}grpSp")
    assert rebuilt_group is not None
    child_ids = [
        element.get("id")
        for element in rebuilt_group.findall(
            f"./{{{namespace}}}sp/{{{namespace}}}nvSpPr/{{{namespace}}}cNvPr"
        )
    ]
    start = rebuilt_group.find(f".//{{{drawing_namespace}}}stCxn")
    end = rebuilt_group.find(f".//{{{drawing_namespace}}}endCxn")
    assert start is not None and end is not None
    assert [start.get("id"), end.get("id")] == child_ids
    assert child_ids != ["220", "340"]


def test_authored_group_emits_picture_child_with_its_media_relationship() -> None:
    image = io.BytesIO()
    Image.new("RGB", (16, 16), "#2563eb").save(image, "PNG")
    group = GroupNode(
        box=Box(x=2_000_000, y=1_000_000, width=3_200_000, height=2_100_000),
        child_box=Box(x=0, y=0, width=1_600_000, height=1_050_000),
        children=(
            ShapeNode(
                box=Box(x=0, y=0, width=800_000, height=500_000),
                fill=PictureFill(data=image.getvalue(), ext="png"),
            ),
        ),
    )

    pptx = build_pptx([SlideIR(width=12_192_000, height=6_858_000, contents=(group,))], faces=[])

    slide_xml = _slide_xml(pptx)
    assert "<p:grpSp>" in slide_xml
    assert "<p:pic>" in slide_xml
    assert "r:embed=" in slide_xml
    assert "relationships/image" in _slide_rels(pptx)
    [recovered] = read_pptx(pptx)[0].contents
    assert isinstance(recovered, GroupNode)
    [picture] = recovered.children
    assert isinstance(picture, ShapeNode)
    assert isinstance(picture.fill, PictureFill)
    assert picture.fill.data == image.getvalue()


def test_authored_group_emits_portable_child_fallback_without_failing_the_deck() -> None:
    image = io.BytesIO()
    Image.new("RGBA", (16, 16), (37, 99, 235, 180)).save(image, "PNG")
    child = ShapeNode(
        box=Box(x=0, y=0, width=800_000, height=500_000),
        fill=SolidFill(color=Rgba(r=239, g=68, b=68)),
        portable_fallback=PortableFallback(
            box=Box(x=-100_000, y=-100_000, width=1_000_000, height=700_000),
            picture=PictureFill(data=image.getvalue(), ext="png"),
        ),
    )
    group = GroupNode(
        box=Box(x=2_000_000, y=1_000_000, width=3_200_000, height=2_100_000),
        child_box=Box(x=0, y=0, width=1_600_000, height=1_050_000),
        children=(child,),
    )

    pptx = build_pptx([SlideIR(width=12_192_000, height=6_858_000, contents=(group,))], faces=[])

    slide_xml = _slide_xml(pptx)
    assert "<p:grpSp>" in slide_xml
    assert "<mc:AlternateContent" in slide_xml
    assert "<mc:Fallback><p:pic>" in slide_xml
    assert "relationships/image" in _slide_rels(pptx)


def test_grouped_text_hyperlink_retains_its_slide_relationship() -> None:
    hyperlink = Hyperlink(url="https://example.com/grouped")
    group = GroupNode(
        box=Box(x=2_000_000, y=1_000_000, width=3_200_000, height=2_100_000),
        child_box=Box(x=0, y=0, width=1_600_000, height=1_050_000),
        children=(
            ShapeNode(
                box=Box(x=0, y=0, width=1_600_000, height=1_050_000),
                text=TextBody(
                    paragraphs=(
                        TextParagraph(
                            runs=(
                                TextRun(
                                    text="Grouped link",
                                    font_family="Arial",
                                    size_pt=18,
                                    hyperlink=hyperlink,
                                ),
                            )
                        ),
                    )
                ),
            ),
        ),
    )

    [recovered] = read_pptx(
        build_pptx([SlideIR(width=12_192_000, height=6_858_000, contents=(group,))], faces=[])
    )[0].contents

    assert isinstance(recovered, GroupNode)
    child = recovered.children[0]
    assert isinstance(child, ShapeNode)
    assert child.text is not None
    assert child.text.paragraphs[0].runs[0].hyperlink == hyperlink


def test_pure_picture_fill_emits_native_picture_with_crop() -> None:
    slide = SlideIR(
        width=12_192_000,
        height=6_858_000,
        shapes=(
            ShapeNode(
                box=Box(x=100, y=200, width=300, height=400),
                fill=PictureFill(data=b"png", crop=SrcRect(left=1 / 3, right=1 / 3)),
            ),
        ),
    )

    xml = _slide_xml(build_pptx([slide], faces=[]))

    assert "<p:pic>" in xml
    assert "<p:blipFill>" in xml
    assert '<a:srcRect l="33333" r="33333"/>' in xml


def test_repeated_bitmap_reuses_one_media_relationship() -> None:
    picture = PictureFill(data=b"same-png")
    slide = SlideIR(
        width=12_192_000,
        height=6_858_000,
        shapes=(
            ShapeNode(box=Box(x=0, y=0, width=100, height=100), fill=picture),
            ShapeNode(box=Box(x=100, y=0, width=100, height=100), fill=picture),
        ),
    )

    pptx = build_pptx([slide], faces=[])
    xml = _slide_xml(pptx)
    rels = _slide_rels(pptx)
    with zipfile.ZipFile(io.BytesIO(pptx)) as archive:
        media = [name for name in archive.namelist() if name.startswith("ppt/media/")]

    assert xml.count('r:embed="rId2"') == 2
    assert rels.count("/relationships/image") == 1
    assert media == ["ppt/media/image1.png"]


def test_repeated_bitmap_reuses_one_media_part_across_slides() -> None:
    picture = PictureFill(data=b"same-png")
    slide = SlideIR(
        width=12_192_000,
        height=6_858_000,
        shapes=(ShapeNode(box=Box(x=0, y=0, width=100, height=100), fill=picture),),
    )

    pptx = build_pptx([slide, slide], faces=[])
    with zipfile.ZipFile(io.BytesIO(pptx)) as archive:
        media = [name for name in archive.namelist() if name.startswith("ppt/media/")]

    assert media == ["ppt/media/image1.png"]
    assert 'Target="../media/image1.png"' in _slide_rels(pptx)
    assert 'Target="../media/image1.png"' in _slide_rels(pptx, "ppt/slides/_rels/slide2.xml.rels")


def test_build_pptx_requires_a_slide() -> None:
    with pytest.raises(ValueError, match="at least one slide"):
        build_pptx([])


def test_build_pptx_rejects_mismatched_slide_sizes() -> None:
    a = SlideIR(width=12_192_000, height=6_858_000, shapes=())
    b = SlideIR(width=9_144_000, height=6_858_000, shapes=())  # 4:3 — different width
    with pytest.raises(ValueError, match="share one size"):
        build_pptx([a, b])


def test_build_pptx_rejects_conflicting_preserved_ambient_themes() -> None:
    def slide(theme_data: bytes) -> SlideIR:
        payload = PreservationPayload(
            kind="graphicFrame",
            root_xml=(
                "<p:graphicFrame "
                'xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"/>'
            ),
            ambient_theme=PreservationPart(
                name="ppt/theme/theme1.xml",
                content_type=("application/vnd.openxmlformats-officedocument.theme+xml"),
                data=theme_data,
            ),
        )
        return SlideIR(
            width=12_192_000,
            height=6_858_000,
            contents=(
                PreservedNode(
                    box=Box(x=0, y=0, width=1_000_000, height=1_000_000),
                    payload=payload,
                ),
            ),
        )

    with pytest.raises(ValueError, match="conflicting preserved ambient themes across slides"):
        build_pptx([slide(b"<theme-one/>"), slide(b"<theme-two/>")], faces=[])


def test_run_underline_and_strike_emit_both_attrs() -> None:
    xml = _slide_xml(build_pptx([_decorated_run(underline=True, strike=True)], faces=[]))
    assert 'u="sng"' in xml
    assert 'strike="sngStrike"' in xml


def test_run_caps_uppercase_emits_cap_all_and_keeps_raw_text() -> None:
    # The IR run text is the authored text; PowerPoint applies the cap. We must NOT pre-uppercase.
    ir = SlideIR(
        width=12_192_000,
        height=6_858_000,
        shapes=(
            ShapeNode(
                box=Box(x=0, y=0, width=3_000_000, height=1_000_000),
                text=TextBody(
                    paragraphs=(
                        TextParagraph(
                            runs=(
                                TextRun(text="Hello", font_family="Inter", size_pt=18, caps="all"),
                            )
                        ),
                    )
                ),
            ),
        ),
    )
    xml = _slide_xml(build_pptx([ir], faces=[]))
    assert 'cap="all"' in xml
    assert "<a:t>Hello</a:t>" in xml  # raw text, not "HELLO"
    assert "HELLO" not in xml


def test_run_small_caps_and_letter_spacing_emit_cap_small_and_spc() -> None:
    xml = _slide_xml(build_pptx([_decorated_run(caps="small", letter_spacing_pt=2.0)], faces=[]))
    assert 'cap="small"' in xml
    assert 'spc="200"' in xml  # 2.0pt → 200 (1/100 pt)


def test_external_hyperlink_emits_hlinkclick_and_external_rel() -> None:
    pptx = build_pptx([_decorated_run(hyperlink=Hyperlink(url="https://example.com"))], faces=[])
    xml = _slide_xml(pptx)
    rels = _slide_rels(pptx)
    assert "<a:hlinkClick" in xml and 'r:id="' in xml
    assert 'Target="https://example.com" TargetMode="External"' in rels
    assert "/relationships/hyperlink" in rels


def test_out_of_range_slide_jump_is_dropped_with_warning() -> None:
    # A rel targeting a slide part that doesn't exist makes PowerPoint repair the file.
    deck = [_decorated_run(hyperlink=Hyperlink(slide_index=1))]  # only 1 slide in deck
    with pytest.warns(UserWarning, match="targets slide 2"):
        pptx = build_pptx(deck, faces=[])
    assert 'Target="slide2.xml"' not in _slide_rels(pptx)
    assert "hlinksldjump" not in _slide_xml(pptx)


def test_slide_jump_hyperlink_emits_jump_action_and_internal_slide_rel() -> None:
    deck = [
        _decorated_run(hyperlink=Hyperlink(slide_index=1)),
        SlideIR(width=12_192_000, height=6_858_000, shapes=()),
    ]
    pptx = build_pptx(deck, faces=[])
    xml = _slide_xml(pptx)
    rels = _slide_rels(pptx)
    assert 'action="ppaction://hlinksldjump"' in xml
    assert 'Target="slide2.xml"' in rels
    assert "/relationships/slide" in rels


# --------------------------------------------------------------------------- paragraph pPr tests


def _para_ir(paragraph: TextParagraph) -> SlideIR:
    """Helper: wrap a single paragraph in a minimal SlideIR."""
    return SlideIR(
        width=12_192_000,
        height=6_858_000,
        shapes=(
            ShapeNode(
                box=Box(x=0, y=0, width=3_000_000, height=1_000_000),
                text=TextBody(paragraphs=(paragraph,)),
            ),
        ),
    )


def test_ppr_line_spacing_percent_emits_spc_pct() -> None:
    """line_spacing(percent=1.5) → <a:lnSpc><a:spcPct val="150000"/></a:lnSpc>."""
    para = TextParagraph(
        runs=(TextRun(text="x", font_family="Arial", size_pt=12),),
        line_spacing=LineSpacing(percent=1.5),
    )
    xml = _slide_xml(build_pptx([_para_ir(para)], faces=[]))
    assert '<a:lnSpc><a:spcPct val="150000"/></a:lnSpc>' in xml


def test_ppr_line_spacing_points_emits_spc_pts() -> None:
    """line_spacing(points=18.0) → <a:lnSpc><a:spcPts val="1800"/></a:lnSpc>."""
    para = TextParagraph(
        runs=(TextRun(text="x", font_family="Arial", size_pt=12),),
        line_spacing=LineSpacing(points=18.0),
    )
    xml = _slide_xml(build_pptx([_para_ir(para)], faces=[]))
    assert '<a:lnSpc><a:spcPts val="1800"/></a:lnSpc>' in xml


def test_ppr_space_before_after_emits_spc_bef_aft() -> None:
    """space_before_pt=9 → <a:spcBef><a:spcPts val="900"/>; space_after_pt=18 → val="1800"."""
    para = TextParagraph(
        runs=(TextRun(text="x", font_family="Arial", size_pt=12),),
        space_before_pt=9.0,
        space_after_pt=18.0,
    )
    xml = _slide_xml(build_pptx([_para_ir(para)], faces=[]))
    assert '<a:spcBef><a:spcPts val="900"/></a:spcBef>' in xml
    assert '<a:spcAft><a:spcPts val="1800"/></a:spcAft>' in xml


def test_ppr_mar_l_and_indent_emit_emu_attrs() -> None:
    """left_margin_pt=36 → marL="457200"; indent_pt=18 → indent="228600" (1pt=12700 EMU)."""
    para = TextParagraph(
        runs=(TextRun(text="x", font_family="Arial", size_pt=12),),
        left_margin_pt=36.0,  # 36 * 12700 = 457200
        indent_pt=18.0,  # 18 * 12700 = 228600
    )
    xml = _slide_xml(build_pptx([_para_ir(para)], faces=[]))
    assert 'marL="457200"' in xml
    assert 'indent="228600"' in xml


def test_ppr_bu_char_emits_buchar_element() -> None:
    """CharBullet(char='•') → <a:buChar char="•"/>."""
    para = TextParagraph(
        runs=(TextRun(text="item", font_family="Arial", size_pt=12),),
        bullet=CharBullet(char="•"),
    )
    xml = _slide_xml(build_pptx([_para_ir(para)], faces=[]))
    assert '<a:buChar char="&#x2022;"/>' in xml or '<a:buChar char="•"/>' in xml


def test_ppr_bu_autonum_emits_buautonum_element() -> None:
    """AutoNumberBullet(scheme='arabicPeriod') → <a:buAutoNum type="arabicPeriod" startAt="1"/>."""
    para = TextParagraph(
        runs=(TextRun(text="item", font_family="Arial", size_pt=12),),
        bullet=AutoNumberBullet(scheme="arabicPeriod"),
    )
    xml = _slide_xml(build_pptx([_para_ir(para)], faces=[]))
    assert 'type="arabicPeriod"' in xml
    assert "a:buAutoNum" in xml


def test_ppr_child_order_lnspc_before_spcbef_before_buchar() -> None:
    """ECMA-376 child order: lnSpc < spcBef < spcAft < buChar within a:pPr."""
    para = TextParagraph(
        runs=(TextRun(text="item", font_family="Arial", size_pt=12),),
        line_spacing=LineSpacing(percent=1.2),
        space_before_pt=6.0,
        space_after_pt=3.0,
        bullet=CharBullet(char="•"),
    )
    xml = _slide_xml(build_pptx([_para_ir(para)], faces=[]))
    lnspc_pos = xml.find("<a:lnSpc>")
    spcbef_pos = xml.find("<a:spcBef>")
    spcaft_pos = xml.find("<a:spcAft>")
    buchar_pos = xml.find("<a:buChar")
    assert lnspc_pos < spcbef_pos < spcaft_pos < buchar_pos, (
        "ECMA child order violated: lnSpc must precede spcBef, spcAft, then buChar"
    )


def test_ppr_lvl_attr_emitted_for_nonzero_level() -> None:
    """level=2 → lvl="2" in a:pPr; level=0 → no lvl attr."""
    para_lvl2 = TextParagraph(
        runs=(TextRun(text="deep", font_family="Arial", size_pt=12),),
        level=2,
        bullet=CharBullet(char="•"),
    )
    para_lvl0 = TextParagraph(
        runs=(TextRun(text="top", font_family="Arial", size_pt=12),),
        level=0,
    )
    xml_lvl2 = _slide_xml(build_pptx([_para_ir(para_lvl2)], faces=[]))
    xml_lvl0 = _slide_xml(build_pptx([_para_ir(para_lvl0)], faces=[]))
    assert 'lvl="2"' in xml_lvl2
    assert 'lvl="0"' not in xml_lvl0
